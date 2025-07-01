import os
import json
from typing import Dict, List, Any, Optional
from pptx import Presentation
from openai import OpenAI
import streamlit as st

class PowerPointAnalyzer:
    def __init__(self, file_content):
        """Initialize the PowerPoint analyzer with file content."""
        self.presentation = Presentation(file_content)
        self.slides_data = []
        self.openai_client = self._initialize_openai()
        
    def _initialize_openai(self):
        """Initialize OpenAI client with API key."""
        api_key = os.environ.get("OPENAI_API_KEY", "dummy-key")
        if api_key == "dummy-key":
            return None
        return OpenAI(api_key=api_key)
    
    def extract_slide_content(self) -> List[Dict[str, Any]]:
        """Extract text content and metadata from all slides."""
        slides_data = []
        
        for i, slide in enumerate(self.presentation.slides):
            slide_data = {
                'slide_number': i + 1,
                'title': '',
                'content': [],
                'bullet_points': [],
                'text_length': 0,
                'shape_count': len(slide.shapes)
            }
            
            # Extract text from all shapes
            for shape in slide.shapes:
                try:
                    # Check if shape has text attribute and text_frame
                    if not hasattr(shape, "text"):
                        continue
                    
                    text = str(getattr(shape, "text", "")).strip()
                    if not text:
                        continue
                    
                    # Try to identify title (usually the first or largest text)
                    if not slide_data['title'] and len(text) < 100:
                        slide_data['title'] = text
                    else:
                        slide_data['content'].append(text)
                    
                    # Check for bullet points
                    if text.startswith(('•', '-', '*')) or '\n•' in text or '\n-' in text:
                        bullet_points = [line.strip() for line in text.split('\n') if line.strip()]
                        slide_data['bullet_points'].extend(bullet_points)
                        
                except (AttributeError, Exception):
                    # Skip shapes that don't have text or cause errors
                    continue
            
            # Calculate total text length
            all_text = slide_data['title'] + ' ' + ' '.join(slide_data['content'])
            slide_data['text_length'] = len(all_text.strip())
            
            slides_data.append(slide_data)
        
        self.slides_data = slides_data
        return slides_data
    
    def get_presentation_overview(self) -> Dict[str, Any]:
        """Get basic overview of the presentation."""
        if not self.slides_data:
            self.extract_slide_content()
        
        total_slides = len(self.slides_data)
        total_text_length = sum(slide['text_length'] for slide in self.slides_data)
        slides_with_titles = len([slide for slide in self.slides_data if slide['title']])
        avg_text_per_slide = total_text_length / total_slides if total_slides > 0 else 0
        
        return {
            'total_slides': total_slides,
            'total_text_length': total_text_length,
            'slides_with_titles': slides_with_titles,
            'average_text_per_slide': round(avg_text_per_slide, 1),
            'slides_with_content': len([slide for slide in self.slides_data if slide['content']]),
            'total_bullet_points': sum(len(slide['bullet_points']) for slide in self.slides_data)
        }
    
    def analyze_content_structure(self) -> Dict[str, Any]:
        """Analyze the structure and organization of content."""
        if not self.slides_data:
            self.extract_slide_content()
        
        # Analyze text distribution
        text_lengths = [slide['text_length'] for slide in self.slides_data]
        
        structure_analysis = {
            'text_distribution': {
                'min_text_length': min(text_lengths) if text_lengths else 0,
                'max_text_length': max(text_lengths) if text_lengths else 0,
                'avg_text_length': sum(text_lengths) / len(text_lengths) if text_lengths else 0
            },
            'content_consistency': self._analyze_consistency(),
            'slide_types': self._categorize_slides()
        }
        
        return structure_analysis
    
    def _analyze_consistency(self) -> Dict[str, Any]:
        """Analyze consistency across slides."""
        if not self.slides_data:
            return {}
        
        # Check for consistent titling
        titled_slides = len([slide for slide in self.slides_data if slide['title']])
        title_consistency = titled_slides / len(self.slides_data) if self.slides_data else 0
        
        # Check text length variation
        text_lengths = [slide['text_length'] for slide in self.slides_data if slide['text_length'] > 0]
        if len(text_lengths) > 1:
            avg_length = sum(text_lengths) / len(text_lengths)
            variance = sum((x - avg_length) ** 2 for x in text_lengths) / len(text_lengths)
            std_dev = variance ** 0.5
            consistency_score = max(0, 1 - (std_dev / avg_length)) if avg_length > 0 else 0
        else:
            consistency_score = 1.0
        
        return {
            'title_consistency': round(title_consistency, 2),
            'content_consistency_score': round(consistency_score, 2),
            'consistent_structure': title_consistency > 0.8 and consistency_score > 0.6
        }
    
    def _categorize_slides(self) -> Dict[str, int]:
        """Categorize slides by type based on content."""
        categories = {
            'title_slides': 0,
            'content_slides': 0,
            'bullet_slides': 0,
            'minimal_slides': 0
        }
        
        for slide in self.slides_data:
            if slide['text_length'] < 50:
                categories['minimal_slides'] += 1
            elif len(slide['bullet_points']) > 3:
                categories['bullet_slides'] += 1
            elif slide['title'] and slide['text_length'] < 200:
                categories['title_slides'] += 1
            else:
                categories['content_slides'] += 1
        
        return categories
    
    def generate_ai_insights(self) -> Dict[str, Any]:
        """Generate AI-powered insights about the presentation."""
        if not self.openai_client:
            return {
                'summary': 'AI analysis not available - API key not configured',
                'key_themes': ['AI analysis requires OpenAI API key'],
                'recommendations': ['Configure OpenAI API key for detailed insights'],
                'content_quality': 'Unable to assess without AI analysis'
            }
        
        if not self.slides_data:
            self.extract_slide_content()
        
        # Prepare content for analysis
        presentation_text = self._prepare_text_for_analysis()
        
        try:
            # Analyze presentation content
            analysis_prompt = f"""
            Analyze this PowerPoint presentation content and provide insights in JSON format:
            
            {presentation_text}
            
            Please provide:
            1. A brief summary of the main topic/theme
            2. Key themes or topics covered (list of 3-5 items)
            3. Content quality assessment (score 1-10 with explanation)
            4. 3-5 specific recommendations for improvement
            
            Respond in JSON format with keys: summary, key_themes, content_quality, recommendations
            """
            
            # the newest OpenAI model is "gpt-4o" which was released May 13, 2024.
            # do not change this unless explicitly requested by the user
            response = self.openai_client.chat.completions.create(
                model="gpt-4o",
                messages=[
                    {
                        "role": "system",
                        "content": "You are an expert presentation analyst. Analyze presentations and provide constructive feedback in JSON format."
                    },
                    {
                        "role": "user",
                        "content": analysis_prompt
                    }
                ],
                response_format={"type": "json_object"}
            )
            
            content = response.choices[0].message.content
            if content:
                return json.loads(content)
            else:
                return {
                    'summary': 'No content returned from AI analysis',
                    'key_themes': ['Analysis incomplete'],
                    'recommendations': ['Retry analysis'],
                    'content_quality': 'Unable to assess'
                }
            
        except Exception as e:
            st.error(f"Error generating AI insights: {str(e)}")
            return {
                'summary': 'Error generating AI analysis',
                'key_themes': ['Analysis failed'],
                'recommendations': ['Check API configuration'],
                'content_quality': 'Unable to assess due to error'
            }
    
    def _prepare_text_for_analysis(self) -> str:
        """Prepare presentation text for AI analysis."""
        content_parts = []
        
        for slide in self.slides_data:
            slide_text = f"Slide {slide['slide_number']}:\n"
            if slide['title']:
                slide_text += f"Title: {slide['title']}\n"
            if slide['content']:
                slide_text += f"Content: {' '.join(slide['content'])}\n"
            content_parts.append(slide_text)
        
        return '\n'.join(content_parts)
    
    def generate_presentation_kpis(self) -> List[Dict[str, Any]]:
        """Generate KPIs specific to presentation analysis."""
        if not self.slides_data:
            self.extract_slide_content()
        
        overview = self.get_presentation_overview()
        structure = self.analyze_content_structure()
        
        kpis = []
        
        # Presentation Length KPI
        if overview['total_slides'] > 0:
            kpis.append({
                'name': 'Presentation Length',
                'value': overview['total_slides'],
                'unit': 'slides',
                'category': 'Structure',
                'description': 'Total number of slides in the presentation',
                'recommendation': self._get_length_recommendation(overview['total_slides'])
            })
        
        # Content Density KPI
        if overview['average_text_per_slide'] > 0:
            kpis.append({
                'name': 'Content Density',
                'value': round(overview['average_text_per_slide'], 1),
                'unit': 'characters/slide',
                'category': 'Content Quality',
                'description': 'Average amount of text content per slide',
                'recommendation': self._get_density_recommendation(overview['average_text_per_slide'])
            })
        
        # Structure Consistency KPI
        consistency_score = structure['content_consistency']['content_consistency_score']
        if consistency_score >= 0:
            kpis.append({
                'name': 'Content Consistency',
                'value': round(consistency_score * 100, 1),
                'unit': '%',
                'category': 'Structure',
                'description': 'How consistent the content length is across slides',
                'recommendation': self._get_consistency_recommendation(consistency_score)
            })
        
        # Title Coverage KPI
        title_coverage = (overview['slides_with_titles'] / overview['total_slides']) * 100 if overview['total_slides'] > 0 else 0
        kpis.append({
            'name': 'Title Coverage',
            'value': round(title_coverage, 1),
            'unit': '%',
            'category': 'Structure',
            'description': 'Percentage of slides with clear titles',
            'recommendation': self._get_title_recommendation(title_coverage)
        })
        
        return kpis
    
    def _get_length_recommendation(self, slide_count: int) -> str:
        """Get recommendation based on slide count."""
        if slide_count < 5:
            return "Consider adding more content or combining with other topics"
        elif slide_count > 20:
            return "Consider breaking into multiple presentations for better audience engagement"
        else:
            return "Good presentation length for audience attention"
    
    def _get_density_recommendation(self, avg_chars: float) -> str:
        """Get recommendation based on content density."""
        if avg_chars < 100:
            return "Consider adding more detailed content to slides"
        elif avg_chars > 500:
            return "Consider reducing text density for better readability"
        else:
            return "Good balance of content per slide"
    
    def _get_consistency_recommendation(self, score: float) -> str:
        """Get recommendation based on consistency score."""
        if score < 0.4:
            return "Improve consistency by standardizing slide content length"
        elif score < 0.7:
            return "Good consistency, minor adjustments could help"
        else:
            return "Excellent content consistency across slides"
    
    def _get_title_recommendation(self, coverage: float) -> str:
        """Get recommendation based on title coverage."""
        if coverage < 50:
            return "Add clear titles to more slides for better navigation"
        elif coverage < 80:
            return "Good title usage, consider adding titles to remaining slides"
        else:
            return "Excellent title coverage for clear presentation structure"